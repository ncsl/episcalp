from pptx import Presentation
from pptx.util import Inches, Pt


def create_title_slide(ppt_fpath, title="", subtitle=""):
    """
    Create a ppt presentation by generating a title slide and saving it.

    Parameters
    ----------
    ppt_fpath: pathlib.Path
        Absolute path of location to save the ppt. Must be a pathlib Path object, not a str
    title: str
        The text you want to be in the title box. Default is empty
    subtitle: str
        The text you want to be in the subtitle box. Default is empty

    Returns
    -------

    """

    """if not ppt_fpath.endswith(".pptx"):
        ppt_fpath += ".pptx"""""  # This does not work on Windows, uncomment if on other system and comment
    # out section below

    ppt_fpath.touch()  # some of the windows functions require the file to exist
    if ".pptx" not in str(ppt_fpath):
        ppt_fpath.unlink()  # Delete the temp file
        ppt_fpath = ppt_fpath.parent / (ppt_fpath.name + ".pptx")  # Rename
        ppt_fpath.touch()  # Create the new file

    # Create a presentation object. The default size for this package is not a usual slide size, so you have to modify
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    # There are set slide layouts. The title slide is not all that important, so we just choose the default one
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)

    # Access the title textbox and add your text
    title_txtbox = slide.shapes.title
    title_txtbox.text = title

    # Access the subtitle textbox and add your text
    subtitle_txtbox = slide.placeholders[1]
    subtitle_txtbox.text = subtitle

    prs.save(ppt_fpath)


def append_slide_to_presentation(ppt_fpath, distribution_fpath, topo_fpath, movie_fpath, subject, info_dict=None):
    """
    Append a slide to an existing presentation.

    This function is just an example where images for a distribution, topographical (static) map, topographical (moving)
    map and basic info are added to a ppt slide.

    The ppt must already exist for this function to work. See `create_title_slide`.

    Parameters
    ----------
    ppt_fpath: pathlib.Path
        Absolute path of location to save the ppt. Must be the same as used in create the presentation.
    distribution_fpath: pathlib.Path
        Absolute path of the distribution plot. Can be any image format.
    topo_fpath: pathlib.Path
        Absolute path of the topographical map. Can be any image format
    movie_fpath: pathlib.Path
        Absolute path of the topographical movie. Only tested with mp4 so far
    subject: str
        Name of the subject to include on the slide
    info_dict: dict
        Dict of extra information about the subject/run you would like to include on the slide.

    Returns
    -------

    """
    # Open the existing presentation
    prs = Presentation(ppt_fpath)
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    # Append a new blank slide
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    # Add distribution plot. Mess around with the size and position parameters on a single slide before automating this
    # in a loop
    left = Inches(10.5)
    top = Inches(0)
    width = Inches(5)
    distribution_img = slide.shapes.add_picture(str(distribution_fpath), left, top, width=width)

    # Add the topomap
    left = Inches(1.5)
    top = Inches(1.5)
    width = Inches(9)
    topomap_img = slide.shapes.add_picture(str(topo_fpath), left, top, width=width)

    # Add the text box that contains subject information. You will have to adjust the width and height parameters
    # depending on how much info you want to include in the slide
    left = top = Inches(0)
    width = Inches(2.2)
    height = Inches(1)
    metadata = slide.shapes.add_textbox(left, top, width, height)
    tf = metadata.text_frame
    subject_txt = tf.add_paragraph()
    subject_txt.text = f"Subject: {subject}"
    subject_txt.font.size = Pt(16)
    for key, value in info_dict.items():
        info_txt = tf.add_paragraph()
        info_txt.text = f"{key}: {value}"
        info_txt.font.size = Pt(16)

    # Add a movie on its own slide
    if movie_fpath:
        left = Inches(4)
        width = Inches(9)
        movie_slide = prs.slides.add_slide(blank_slide_layout)
        movie_slide.shapes.add_movie(movie_fpath, left, top, width=width, height=width)

    prs.save(ppt_fpath)
